from __future__ import annotations

import argparse
import re
from pathlib import Path
from typing import Any, Iterable

import yaml
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement


DEFAULT_TEMPLATE = Path(__file__).with_name("RS 세미나 템플릿.pptx")
DEFAULT_OUTPUT_DIR = Path(__file__).with_name("results")
DEFAULT_EOD_TITLE = "EOD"


class ConfigError(ValueError):
    """Raised when the YAML input or PPTX generation setup is invalid."""


# Tags related to bullet formatting in DrawingML paragraph properties.
# A body placeholder can inherit bullet formatting from the PowerPoint layout.
# We explicitly replace inherited bullet settings with <a:buNone/> for plain body text.
BULLET_RELATED_TAGS = {
    qn("a:buClrTx"),
    qn("a:buClr"),
    qn("a:buSzTx"),
    qn("a:buSzPct"),
    qn("a:buSzPts"),
    qn("a:buFontTx"),
    qn("a:buFont"),
    qn("a:buNone"),
    qn("a:buAutoNum"),
    qn("a:buChar"),
    qn("a:buBlip"),
}

# XML child order matters. Bullet elements should appear before these elements if present.
BU_NONE_INSERT_BEFORE_TAGS = {
    qn("a:tabLst"),
    qn("a:defRPr"),
    qn("a:extLst"),
}


def slugify(value: str, fallback: str = "rs_seminar") -> str:
    value = value.strip().lower()
    value = re.sub(r"[^0-9a-zA-Z가-힣._-]+", "_", value)
    value = re.sub(r"_+", "_", value).strip("._-")
    return value or fallback


def load_yaml(path: Path) -> dict[str, Any]:
    if not path.exists():
        raise FileNotFoundError(f"YAML file not found: {path}")
    with path.open("r", encoding="utf-8") as f:
        data = yaml.safe_load(f)
    if not isinstance(data, dict):
        raise ConfigError("YAML root must be a mapping/object.")
    return data


def as_required_mapping(data: dict[str, Any], key: str) -> dict[str, Any]:
    value = data.get(key)
    if not isinstance(value, dict):
        raise ConfigError(f"'{key}' must be a mapping/object.")
    return value


def as_required_list(data: dict[str, Any], key: str) -> list[Any]:
    value = data.get(key)
    if not isinstance(value, list):
        raise ConfigError(f"'{key}' must be a list.")
    return value


def required_str(data: dict[str, Any], key: str, context: str) -> str:
    value = data.get(key)
    if value is None:
        raise ConfigError(f"Missing required field: {context}.{key}")
    if not isinstance(value, str):
        raise ConfigError(f"{context}.{key} must be a string.")
    return value


def optional_bool(data: dict[str, Any], key: str, default: bool, context: str) -> bool:
    value = data.get(key, default)
    if not isinstance(value, bool):
        raise ConfigError(f"{context}.{key} must be a boolean when provided.")
    return value


def normalize_deck(data: dict[str, Any]) -> dict[str, Any]:
    deck = data.get("deck") or {}
    if not isinstance(deck, dict):
        raise ConfigError("'deck' must be a mapping/object when provided.")

    # Default is deliberately False: all body content is rendered as plain paragraphs.
    # Set deck.enable_bullets: true only if explicit PowerPoint bullets are desired.
    deck = dict(deck)
    deck["enable_bullets"] = optional_bool(deck, "enable_bullets", False, "deck")
    return deck


def normalize_cover(data: dict[str, Any]) -> dict[str, str]:
    cover = as_required_mapping(data, "cover")
    return {
        "title": required_str(cover, "title", "cover"),
        "subtitle": required_str(cover, "subtitle", "cover"),
        "date": required_str(cover, "date", "cover"),
        "venue": required_str(cover, "venue", "cover"),
    }


def normalize_main(data: dict[str, Any]) -> dict[str, dict[str, Any]]:
    raw_main = as_required_list(data, "main")
    result: dict[str, dict[str, Any]] = {}
    for i, item in enumerate(raw_main):
        context = f"main[{i}]"
        if not isinstance(item, dict):
            raise ConfigError(f"{context} must be a mapping/object.")
        slide_id = required_str(item, "id", context)
        if slide_id in result:
            raise ConfigError(f"Duplicate main slide id: {slide_id}")
        title = required_str(item, "title", context)
        if "body" not in item:
            raise ConfigError(f"Missing required field: {context}.body")
        result[slide_id] = {"id": slide_id, "title": title, "body": item["body"]}
    return result


def normalize_sections(data: dict[str, Any], main_by_id: dict[str, dict[str, Any]]) -> list[dict[str, Any]]:
    raw_sections = as_required_list(data, "sections")
    result: list[dict[str, Any]] = []
    seen_ids: set[str] = set()
    referenced_main_ids: set[str] = set()

    for i, item in enumerate(raw_sections):
        context = f"sections[{i}]"
        if not isinstance(item, dict):
            raise ConfigError(f"{context} must be a mapping/object.")

        section_id = required_str(item, "id", context)
        if section_id in seen_ids:
            raise ConfigError(f"Duplicate section id: {section_id}")
        seen_ids.add(section_id)

        title = required_str(item, "title", context)
        contents = item.get("contents")
        if not isinstance(contents, list):
            raise ConfigError(f"{context}.contents must be a list of main slide ids.")

        normalized_contents: list[str] = []
        for j, slide_id in enumerate(contents):
            if not isinstance(slide_id, str):
                raise ConfigError(f"{context}.contents[{j}] must be a string.")
            if slide_id not in main_by_id:
                raise ConfigError(f"{context}.contents[{j}] references unknown main id: {slide_id}")
            normalized_contents.append(slide_id)
            referenced_main_ids.add(slide_id)

        result.append({"id": section_id, "title": title, "contents": normalized_contents})

    unused_ids = set(main_by_id) - referenced_main_ids
    if unused_ids:
        raise ConfigError("main slides not referenced by any section: " + ", ".join(sorted(unused_ids)))
    return result


def normalize_config(data: dict[str, Any]) -> dict[str, Any]:
    deck = normalize_deck(data)
    cover = normalize_cover(data)
    main_by_id = normalize_main(data)
    sections = normalize_sections(data, main_by_id)
    return {
        "schema_version": data.get("schema_version", "1.0"),
        "deck": deck,
        "cover": cover,
        "sections": sections,
        "main_by_id": main_by_id,
    }


def require_template_slides(prs: Presentation) -> None:
    if len(prs.slides) < 3:
        raise ConfigError(
            "Template must contain at least 3 slides: "
            "slide 1=cover, slide 2=section, slide 3=main."
        )


def delete_all_slides(prs: Presentation) -> None:
    """Remove all sample slides from the template.

    python-pptx has no public delete-slide API, so this uses the underlying
    presentation relationship list. Keep this function small and isolated.
    """
    slide_id_list = prs.slides._sldIdLst  # pylint: disable=protected-access
    for slide_id in list(slide_id_list):
        prs.part.drop_rel(slide_id.rId)
        slide_id_list.remove(slide_id)


def iter_text_shapes(slide: Any) -> Iterable[Any]:
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            yield shape


def iter_placeholders(slide: Any) -> Iterable[Any]:
    for shape in slide.shapes:
        if getattr(shape, "is_placeholder", False):
            yield shape


def placeholders_by_type(slide: Any, placeholder_type: PP_PLACEHOLDER) -> list[Any]:
    matches: list[Any] = []
    for shape in iter_placeholders(slide):
        try:
            if shape.placeholder_format.type == placeholder_type:
                matches.append(shape)
        except ValueError:
            continue
    return matches


def find_title_shape(slide: Any) -> Any:
    if slide.shapes.title is not None:
        return slide.shapes.title

    title_candidates = placeholders_by_type(slide, PP_PLACEHOLDER.TITLE)
    title_candidates += placeholders_by_type(slide, PP_PLACEHOLDER.CENTER_TITLE)
    if title_candidates:
        return sorted(title_candidates, key=lambda s: (s.top, s.left))[0]

    text_shapes = list(iter_text_shapes(slide))
    if not text_shapes:
        raise ConfigError("No title/text shape found on slide.")
    return sorted(text_shapes, key=lambda s: (s.top, s.left))[0]


def find_body_shape(slide: Any) -> Any:
    body_candidates = placeholders_by_type(slide, PP_PLACEHOLDER.BODY)
    if body_candidates:
        return sorted(body_candidates, key=lambda s: (s.top, s.left))[0]

    title_shape = find_title_shape(slide)
    text_shapes = [s for s in iter_text_shapes(slide) if s != title_shape]
    if not text_shapes:
        raise ConfigError("No body placeholder/text shape found on main slide.")

    # Body is normally the largest non-title text region.
    return sorted(text_shapes, key=lambda s: (s.width * s.height), reverse=True)[0]


def force_no_bullet(paragraph: Any) -> None:
    """Force a paragraph to be non-bulleted.

    PowerPoint body placeholders often inherit bullet formatting from the layout.
    python-pptx does not expose a stable public paragraph.bullet = False API, so
    this writes the DrawingML paragraph property <a:buNone/> directly.
    """
    p_pr = paragraph._p.get_or_add_pPr()  # pylint: disable=protected-access

    for child in list(p_pr):
        if child.tag in BULLET_RELATED_TAGS:
            p_pr.remove(child)

    bu_none = OxmlElement("a:buNone")

    insert_idx = len(p_pr)
    for i, child in enumerate(p_pr):
        if child.tag in BU_NONE_INSERT_BEFORE_TAGS:
            insert_idx = i
            break

    p_pr.insert(insert_idx, bu_none)
    paragraph.level = 0


def set_shape_text(shape: Any, text: str, *, no_bullet: bool = True) -> None:
    tf = shape.text_frame
    tf.clear()
    paragraph = tf.paragraphs[0]
    paragraph.text = text or ""
    if no_bullet:
        force_no_bullet(paragraph)


def add_paragraph(
    tf: Any,
    text: str,
    *,
    level: int = 0,
    first: bool = False,
    no_bullet: bool = True,
) -> None:
    paragraph = tf.paragraphs[0] if first else tf.add_paragraph()
    paragraph.text = text or ""
    paragraph.level = max(0, min(int(level), 8))
    if no_bullet:
        force_no_bullet(paragraph)


def collapse_excess_blank_lines(lines: list[str]) -> list[str]:
    """Keep blank lines, but collapse runs of 2+ blank lines to a single blank line."""
    result: list[str] = []
    previous_blank = False
    for line in lines:
        is_blank = not line.strip()
        if is_blank and previous_blank:
            continue
        result.append(line)
        previous_blank = is_blank
    return result or [""]


def render_paragraph_body(shape: Any, text: str) -> None:
    tf = shape.text_frame
    tf.clear()

    lines = text.splitlines() or [""]
    lines = collapse_excess_blank_lines(lines)

    for i, line in enumerate(lines):
        add_paragraph(tf, line, level=0, first=(i == 0), no_bullet=True)


def normalize_bullet_item(item: Any) -> tuple[str, int]:
    if isinstance(item, str):
        return item, 0
    if isinstance(item, dict):
        text = item.get("text")
        if not isinstance(text, str):
            raise ConfigError("Bullet item object must contain string field 'text'.")
        level = item.get("level", 0)
        if not isinstance(level, int):
            raise ConfigError("Bullet item field 'level' must be an integer.")
        return text, level
    raise ConfigError("Bullet item must be either a string or an object with text/level.")


def render_list_as_plain_body(shape: Any, items: list[Any]) -> None:
    """Render list-style YAML as plain paragraphs, not PowerPoint bullets.

    This is the default behavior because the requested output should put all
    main content into normal body text. Level information is preserved only as
    leading indentation spaces.
    """
    tf = shape.text_frame
    tf.clear()

    normalized = [normalize_bullet_item(item) for item in items]
    if not normalized:
        normalized = [("", 0)]

    for i, (text, level) in enumerate(normalized):
        indent = "  " * max(0, min(int(level), 8))
        add_paragraph(tf, indent + text, level=0, first=(i == 0), no_bullet=True)


def render_bullet_body(shape: Any, items: list[Any]) -> None:
    """Render explicit PowerPoint bullets.

    This path is used only when deck.enable_bullets is true and the YAML body is
    explicitly typed as bullets or provided as a list. The default is plain text.
    """
    tf = shape.text_frame
    tf.clear()

    normalized = [normalize_bullet_item(item) for item in items]
    if not normalized:
        normalized = [("", 0)]

    for i, (text, level) in enumerate(normalized):
        add_paragraph(tf, text, level=level, first=(i == 0), no_bullet=False)


def render_body(shape: Any, body: Any, *, enable_bullets: bool = False) -> None:
    if isinstance(body, str):
        render_paragraph_body(shape, body)
        return

    if isinstance(body, list):
        if enable_bullets:
            render_bullet_body(shape, body)
        else:
            render_list_as_plain_body(shape, body)
        return

    if isinstance(body, dict):
        body_type = body.get("type", "paragraph")

        if body_type == "paragraph":
            text = body.get("text", "")
            if not isinstance(text, str):
                raise ConfigError("body.text must be a string when body.type is 'paragraph'.")
            render_paragraph_body(shape, text)
            return

        if body_type == "bullets":
            items = body.get("items", [])
            if not isinstance(items, list):
                raise ConfigError("body.items must be a list when body.type is 'bullets'.")
            if enable_bullets:
                render_bullet_body(shape, items)
            else:
                render_list_as_plain_body(shape, items)
            return

        raise ConfigError(f"Unsupported body.type: {body_type}")

    raise ConfigError("body must be a string, list, or object.")


def fill_cover(slide: Any, cover: dict[str, str]) -> None:
    set_shape_text(find_title_shape(slide), cover["title"])

    subtitles = placeholders_by_type(slide, PP_PLACEHOLDER.SUBTITLE)
    if len(subtitles) < 3:
        title_shape = find_title_shape(slide)
        subtitles = [s for s in iter_text_shapes(slide) if s != title_shape]

    subtitles = sorted(subtitles, key=lambda s: (s.top, s.left))
    if len(subtitles) < 3:
        raise ConfigError(
            "Cover layout must have at least three subtitle/text placeholders "
            "for date, venue, and subtitle."
        )

    top_y = min(s.top for s in subtitles)

    # Group shapes in the top band. Use a generous threshold because EMU values vary by template.
    top_band = [s for s in subtitles if abs(s.top - top_y) <= 300000]
    if len(top_band) >= 2:
        date_shape = sorted(top_band, key=lambda s: s.left)[0]
        venue_shape = sorted(top_band, key=lambda s: s.left)[-1]
        remaining = [s for s in subtitles if s is not date_shape and s is not venue_shape]
        subtitle_shape = sorted(remaining, key=lambda s: s.top)[-1]
    else:
        date_shape, venue_shape = subtitles[0], subtitles[1]
        subtitle_shape = subtitles[-1]

    set_shape_text(date_shape, cover["date"])
    set_shape_text(venue_shape, cover["venue"])
    set_shape_text(subtitle_shape, cover["subtitle"])


def fill_section(slide: Any, title: str) -> None:
    set_shape_text(find_title_shape(slide), title)


def fill_main(slide: Any, item: dict[str, Any], *, enable_bullets: bool = False) -> None:
    set_shape_text(find_title_shape(slide), item["title"])
    render_body(find_body_shape(slide), item["body"], enable_bullets=enable_bullets)


def build_presentation(config: dict[str, Any], template_path: Path, output_path: Path) -> Path:
    if not template_path.exists():
        raise FileNotFoundError(f"Template file not found: {template_path}")

    prs = Presentation(str(template_path))
    require_template_slides(prs)

    cover_layout = prs.slides[0].slide_layout
    section_layout = prs.slides[1].slide_layout
    main_layout = prs.slides[2].slide_layout

    delete_all_slides(prs)

    enable_bullets = bool(config["deck"].get("enable_bullets", False))

    cover_slide = prs.slides.add_slide(cover_layout)
    fill_cover(cover_slide, config["cover"])

    for section in config["sections"]:
        section_slide = prs.slides.add_slide(section_layout)
        fill_section(section_slide, section["title"])

        for slide_id in section["contents"]:
            main_slide = prs.slides.add_slide(main_layout)
            fill_main(main_slide, config["main_by_id"][slide_id], enable_bullets=enable_bullets)

    eod_title = config["deck"].get("eod_title", DEFAULT_EOD_TITLE)
    if not isinstance(eod_title, str):
        raise ConfigError("deck.eod_title must be a string when provided.")

    eod_slide = prs.slides.add_slide(section_layout)
    fill_section(eod_slide, eod_title)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def resolve_output_path(args: argparse.Namespace, config: dict[str, Any]) -> Path:
    if args.out:
        return Path(args.out)
    deck = config["deck"]
    output_name = deck.get("output_name")
    if not isinstance(output_name, str):
        output_name = "rs_seminar"
    return DEFAULT_OUTPUT_DIR / f"{slugify(output_name)}.pptx"


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Fill the RS seminar PPTX template from YAML.")
    parser.add_argument("--config", type=Path, required=True, help="YAML input file, e.g. config/target_example.yaml")
    parser.add_argument("--template", type=Path, default=DEFAULT_TEMPLATE, help="PPTX template path")
    parser.add_argument("--out", type=Path, help="Output PPTX path")
    return parser.parse_args()


def main() -> None:
    args = parse_args()

    try:
        raw_config = load_yaml(args.config)
        config = normalize_config(raw_config)
        output_path = resolve_output_path(args, config)
        created = build_presentation(config, args.template, output_path)
        print(f"Saved at {created.resolve()}")

    except (ConfigError, FileNotFoundError) as exc:
        raise SystemExit(f"Error: {exc}") from exc


if __name__ == "__main__":
    main()
